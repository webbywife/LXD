"""
pptx_builder.py
PowerPoint generation for SKOOLED-AI lesson plans.
Color scheme: BG #0d1f2d (dark navy), accent #00bbd6 (cyan), secondary #faa32b (orange)
"""

import re
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Brand Colors ─────────────────────────────────────────────
NAVY    = RGBColor(0x0d, 0x1f, 0x2d)
NAVY2   = RGBColor(0x0a, 0x16, 0x22)   # slightly darker for gradient effect
CYAN    = RGBColor(0x00, 0xbb, 0xd6)
ORANGE  = RGBColor(0xfa, 0xa3, 0x2b)
GREEN   = RGBColor(0x00, 0xb8, 0x94)
WHITE   = RGBColor(0xff, 0xff, 0xff)
GREY    = RGBColor(0xaa, 0xcc, 0xd4)
DIMWHITE = RGBColor(0xcc, 0xe8, 0xee)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Content area constants
CONTENT_LEFT  = Inches(0.55)
CONTENT_TOP   = Inches(1.65)
CONTENT_W     = Inches(12.2)
CONTENT_H     = SLIDE_H - Inches(1.65) - Inches(0.45)   # bottom margin


# ── Low-level helpers ─────────────────────────────────────────

def _set_bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


# ── Rich bullet list ──────────────────────────────────────────

def _auto_font_size(n_items):
    """Return font size and line spacing based on item count."""
    if n_items <= 3:
        return 16, 1.5   # font pt, line_spacing multiplier
    elif n_items <= 5:
        return 14, 1.45
    elif n_items <= 7:
        return 13, 1.4
    else:
        return 12, 1.35


def _rich_bullets(slide, items, left=CONTENT_LEFT, top=CONTENT_TOP,
                  width=CONTENT_W, height=CONTENT_H,
                  bullet_char="▸", accent_color=None):
    """
    Single textbox containing all bullet paragraphs with proper
    line spacing, space_before, and auto font sizing.
    """
    if not items:
        return

    fs, ls = _auto_font_size(len(items))
    sp_before = Pt(8) if len(items) <= 5 else Pt(5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        item = item.lstrip("•▸-* ").strip()
        if not item:
            continue
        # Truncate very long items gracefully
        if len(item) > 160:
            item = item[:157] + "…"

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = ls           # relative line spacing (e.g. 1.4 = 140%)
        p.space_before = sp_before
        p.space_after = Pt(2)

        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.size = Pt(fs)
        run.font.color.rgb = WHITE
        run.font.bold = False


def _split_items(items, max_per_slide=5):
    """Split a list into chunks for overflow slides."""
    chunks = []
    for i in range(0, len(items), max_per_slide):
        chunks.append(items[i:i + max_per_slide])
    return chunks if chunks else [[]]


# ── Chrome shared by all content slides ──────────────────────

def _add_header(slide, title, subtitle="", title_color=NAVY):
    """Cyan header bar with title + optional subtitle."""
    bar_h = Inches(1.4) if not subtitle else Inches(1.5)
    _add_rect(slide, 0, 0, SLIDE_W, bar_h, CYAN)
    # Left accent strip
    _add_rect(slide, 0, 0, Inches(0.06), bar_h, ORANGE)

    title_top = Inches(0.14) if not subtitle else Inches(0.1)
    _add_textbox(slide, title,
                 Inches(0.55), title_top, Inches(12.2), Inches(0.78),
                 font_size=30, bold=True, color=title_color)
    if subtitle:
        _add_textbox(slide, subtitle,
                     Inches(0.55), Inches(0.94), Inches(12.2), Inches(0.42),
                     font_size=13, color=RGBColor(0x0a, 0x3a, 0x4e))


def _add_footer(slide, text="SKOOLED-AI  ·  Philippine DepEd MATATAG"):
    """Slim dark footer bar with slide label."""
    bar_top = SLIDE_H - Inches(0.36)
    _add_rect(slide, 0, bar_top, SLIDE_W, Inches(0.36), NAVY2)
    _add_textbox(slide, text,
                 Inches(0.5), bar_top + Inches(0.05), Inches(12.3), Inches(0.26),
                 font_size=9, color=GREY, align=PP_ALIGN.LEFT)


def _content_slide(prs, title, subtitle=""):
    """Create a blank content slide with header + footer, return (slide, content_top)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_header(slide, title, subtitle)
    _add_footer(slide)
    content_top = Inches(1.55) if not subtitle else Inches(1.68)
    return slide, content_top


# ── Slide builders ────────────────────────────────────────────

def _slide_title(prs, parsed):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)

    # Left accent strip
    _add_rect(slide, 0, 0, Inches(0.22), SLIDE_H, CYAN)
    # Bottom accent bar
    _add_rect(slide, 0, SLIDE_H - Inches(0.25), SLIDE_W, Inches(0.25), ORANGE)
    # Subtle mid stripe
    _add_rect(slide, Inches(0.22), Inches(3.5), SLIDE_W - Inches(0.22), Inches(0.015),
              RGBColor(0x00, 0xbb, 0xd6))

    # Brand label
    _add_textbox(slide, "SKOOLED-AI",
                 Inches(0.7), Inches(0.9), Inches(11.6), Inches(0.55),
                 font_size=13, bold=False, color=CYAN, align=PP_ALIGN.CENTER)

    # Main title
    title = parsed["title"] or "Lesson Plan"
    _add_textbox(slide, title,
                 Inches(0.7), Inches(1.55), Inches(11.6), Inches(1.9),
                 font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Meta pills row
    meta_parts = []
    if parsed["subject"]:
        meta_parts.append(parsed["subject"])
    if parsed["grade"]:
        meta_parts.append(parsed["grade"])
    if parsed["quarter"]:
        meta_parts.append(f"Quarter {parsed['quarter']}")
    meta = "     ·     ".join(meta_parts)
    if meta:
        _add_textbox(slide, meta,
                     Inches(0.7), Inches(3.6), Inches(11.6), Inches(0.55),
                     font_size=17, color=GREY, align=PP_ALIGN.CENTER)

    # DepEd label above bottom bar
    _add_textbox(slide, "Philippine DepEd  ·  MATATAG Curriculum  ·  AI-Powered Lesson Plan",
                 Inches(0.7), SLIDE_H - Inches(0.72), Inches(11.6), Inches(0.36),
                 font_size=10, color=DIMWHITE, align=PP_ALIGN.CENTER)


def _slides_bullets(prs, title, items, subtitle="", bullet_color=CYAN,
                    max_per=5, continuation_label="(continued)"):
    """
    Build one or more slides for a bullet list, auto-splitting if needed.
    """
    chunks = _split_items(items, max_per)
    for ci, chunk in enumerate(chunks):
        sub = subtitle if ci == 0 else f"{subtitle}  {continuation_label}" if subtitle else continuation_label
        slide, ct = _content_slide(prs, title, sub)
        content_h = SLIDE_H - ct - Inches(0.42)
        _rich_bullets(slide, chunk, top=ct, height=content_h)


def _slide_skills_chips(prs, parsed):
    """21st Century Skills as large coloured chips."""
    items = parsed["skills"] or ["Critical Thinking", "Communication", "Collaboration"]
    slide, _ = _content_slide(prs, "21st Century Skills", "Integrated competencies for this lesson")

    colors = [CYAN, ORANGE, GREEN,
              RGBColor(0x9b, 0x59, 0xb6),
              RGBColor(0x34, 0x98, 0xdb),
              RGBColor(0xe7, 0x4c, 0x3c)]

    # 2-column chip layout
    chip_w = Inches(5.7)
    chip_h = Inches(0.72)
    gap_x  = Inches(0.5)
    gap_y  = Inches(0.18)
    x0     = Inches(0.55)
    y0     = Inches(1.72)

    for i, skill in enumerate(items[:8]):
        col = i % 2
        row = i // 2
        x = x0 + col * (chip_w + gap_x)
        y = y0 + row * (chip_h + gap_y)

        color = colors[i % len(colors)]
        _add_rect(slide, x, y, chip_w, chip_h, color)
        # Left stripe inside chip
        _add_rect(slide, x, y, Inches(0.06), chip_h, NAVY)

        skill_text = skill.lstrip("•-* ").strip()
        if len(skill_text) > 60:
            skill_text = skill_text[:57] + "…"
        _add_textbox(slide, skill_text,
                     x + Inches(0.18), y + Inches(0.14),
                     chip_w - Inches(0.24), chip_h - Inches(0.18),
                     font_size=13, bold=True, color=NAVY)


def _slide_materials_2col(prs, parsed):
    """Materials in a clean 2-column layout."""
    items = parsed["materials"] or ["Textbook", "Whiteboard", "Digital projector"]
    slide, ct = _content_slide(prs, "Materials & Technology")

    mid = (len(items) + 1) // 2
    left_items  = items[:mid]
    right_items = items[mid:]

    col_w = Inches(5.8)
    col_h = SLIDE_H - ct - Inches(0.42)

    _rich_bullets(slide, left_items,
                  left=Inches(0.55), top=ct, width=col_w, height=col_h)
    _rich_bullets(slide, right_items,
                  left=Inches(6.85), top=ct, width=col_w, height=col_h)


def _slide_phase_split(prs, phase):
    """
    Lesson phase slide.  If content has 'teacher:' / 'student:' keywords,
    split into two columns; otherwise single-column with accent label.
    """
    content = phase["content"] or ["Teacher facilitates this phase of the lesson."]
    name    = phase["name"]

    # Detect teacher/student split
    teacher_items, student_items, other = [], [], []
    for item in content:
        ll = item.lower()
        if ll.startswith(("teacher", "t:")):
            teacher_items.append(re.sub(r'^teacher\s*:\s*', '', item, flags=re.I).strip())
        elif ll.startswith(("student", "learner", "s:")):
            student_items.append(re.sub(r'^(student|learner)s?\s*:\s*', '', item, flags=re.I).strip())
        else:
            other.append(item)

    if teacher_items and student_items:
        # Two-column layout
        slide, ct = _content_slide(prs, name)
        col_w = Inches(5.8)
        col_h = SLIDE_H - ct - Inches(0.42)

        # Teacher column header
        _add_rect(slide, Inches(0.55), ct, col_w, Inches(0.35), CYAN)
        _add_textbox(slide, "👩‍🏫 Teacher",
                     Inches(0.65), ct + Inches(0.04), col_w - Inches(0.2), Inches(0.28),
                     font_size=12, bold=True, color=NAVY)
        _rich_bullets(slide, teacher_items,
                      left=Inches(0.55), top=ct + Inches(0.38),
                      width=col_w, height=col_h - Inches(0.38))

        # Student column header
        _add_rect(slide, Inches(6.85), ct, col_w, Inches(0.35), ORANGE)
        _add_textbox(slide, "🧑‍🎓 Students",
                     Inches(6.95), ct + Inches(0.04), col_w - Inches(0.2), Inches(0.28),
                     font_size=12, bold=True, color=NAVY)
        _rich_bullets(slide, student_items,
                      left=Inches(6.85), top=ct + Inches(0.38),
                      width=col_w, height=col_h - Inches(0.38))
    else:
        # Single column — but split if too many items
        _slides_bullets(prs, name, other + teacher_items + student_items, max_per=5)
        return   # already added slides

    return


def _slide_differentiation(prs, parsed):
    """3-column differentiation slide."""
    slide, ct = _content_slide(prs, "Differentiation & Scaffolding")

    diff  = parsed["differentiation"]
    cols  = [
        ("Struggling Learners",      diff.get("struggling") or ["Use visual aids and simplified text."], ORANGE),
        ("Advanced Learners",        diff.get("advanced")   or ["Extend with enrichment tasks."],        CYAN),
        ("ELL / Language Learners",  diff.get("ell")        or ["Provide bilingual vocabulary support."], GREEN),
    ]

    col_w = Inches(3.9)
    gap   = Inches(0.27)
    x0    = Inches(0.55)
    hdr_h = Inches(0.38)
    col_h = SLIDE_H - ct - hdr_h - Inches(0.55)

    for ci, (label, items, color) in enumerate(cols):
        x = x0 + ci * (col_w + gap)

        # Column header
        _add_rect(slide, x, ct, col_w, hdr_h, color)
        _add_textbox(slide, label,
                     x + Inches(0.1), ct + Inches(0.06),
                     col_w - Inches(0.15), hdr_h - Inches(0.08),
                     font_size=12, bold=True, color=NAVY)

        # Bullets inside column
        _rich_bullets(slide, items[:5],
                      left=x + Inches(0.06), top=ct + hdr_h + Inches(0.1),
                      width=col_w - Inches(0.1), height=col_h)


def _slide_assessment(prs, parsed):
    """2-column formative / summative assessment slide."""
    slide, ct = _content_slide(prs, "Assessment")

    assess = parsed["assessment"]
    cols = [
        ("Formative Assessment",  assess.get("formative") or ["Exit ticket", "Observation"], CYAN),
        ("Summative Assessment",  assess.get("summative") or ["Unit test", "Project"],       ORANGE),
    ]

    col_w = Inches(5.8)
    gap   = Inches(0.73)
    x0    = Inches(0.55)
    hdr_h = Inches(0.38)
    col_h = SLIDE_H - ct - hdr_h - Inches(0.55)

    for ci, (label, items, color) in enumerate(cols):
        x = x0 + ci * (col_w + gap)

        _add_rect(slide, x, ct, col_w, hdr_h, color)
        _add_textbox(slide, label,
                     x + Inches(0.1), ct + Inches(0.06),
                     col_w - Inches(0.15), hdr_h - Inches(0.08),
                     font_size=13, bold=True, color=NAVY)

        _rich_bullets(slide, items[:6],
                      left=x + Inches(0.06), top=ct + hdr_h + Inches(0.1),
                      width=col_w - Inches(0.1), height=col_h)


def _slide_reflection(prs, parsed):
    items = parsed["reflection"] or [
        "What went well in this lesson?",
        "What would I change next time?",
        "Did all students meet the learning objectives?",
    ]
    slide, ct = _content_slide(prs, "Teacher Reflection",
                               "Post-lesson self-evaluation prompts")
    col_h = SLIDE_H - ct - Inches(0.42)
    _rich_bullets(slide, items, top=ct, height=col_h, bullet_char="◆")


# ── Markdown Parser ───────────────────────────────────────────

def parse_lesson_markdown(md):
    data = {
        "title": "", "subject": "", "grade": "", "quarter": "",
        "objectives": [], "materials": [], "skills": [], "phases": [],
        "differentiation": {"struggling": [], "advanced": [], "ell": []},
        "assessment": {"formative": [], "summative": []},
        "reflection": [],
    }

    lines = md.split('\n')

    for line in lines:
        if line.startswith('# '):
            data["title"] = line[2:].strip()
            break

    for line in lines[:40]:
        cl = line.strip().lower()
        if 'subject' in cl or 'learning area' in cl:
            m = re.search(r'[:\|]\s*(.+)', line)
            if m: data["subject"] = m.group(1).strip().rstrip('*').strip()
        if 'grade' in cl:
            m = re.search(r'[:\|]\s*(.+)', line)
            if m: data["grade"] = m.group(1).strip().rstrip('*').strip()
        if 'quarter' in cl:
            m = re.search(r'[:\|]\s*(.+)', line)
            if m: data["quarter"] = m.group(1).strip().rstrip('*').strip()

    # Split H2 sections
    sections, current_sec, current_lines = {}, None, []
    for line in lines:
        if line.startswith('## '):
            if current_sec is not None:
                sections[current_sec.lower()] = '\n'.join(current_lines)
            current_sec, current_lines = line[3:].strip(), []
        elif current_sec is not None:
            current_lines.append(line)
    if current_sec:
        sections[current_sec.lower()] = '\n'.join(current_lines)

    def extract_bullets(text):
        items = []
        for line in text.split('\n'):
            line = line.strip()
            if line.startswith(('- ', '* ', '• ')):
                items.append(line[2:].strip())
            elif re.match(r'^\d+\.\s', line):
                items.append(re.sub(r'^\d+\.\s', '', line).strip())
            elif line.startswith('**') and ':' in line:
                items.append(line.replace('**', '').strip())
        return [i for i in items if i]

    for k in sections:
        if 'objective' in k or 'swbat' in k:
            data["objectives"] = extract_bullets(sections[k])[:8]; break
    for k in sections:
        if 'material' in k or 'technolog' in k or 'resource' in k:
            data["materials"] = extract_bullets(sections[k])[:12]; break
    for k in sections:
        if '21st' in k or 'skill' in k:
            data["skills"] = extract_bullets(sections[k])[:8]; break
    for k in sections:
        if 'procedure' in k or 'lesson proper' in k or 'activit' in k:
            sec_text = sections[k]
            parts = re.split(r'\n###\s+', sec_text)
            if len(parts) > 1:
                for pp in parts[1:]:
                    ph_lines = pp.split('\n')
                    ph_name = ph_lines[0].strip()
                    ph_content = extract_bullets('\n'.join(ph_lines[1:]))
                    if not ph_content:
                        ph_content = [l.strip() for l in ph_lines[1:6] if l.strip()]
                    data["phases"].append({"name": ph_name, "content": ph_content[:8]})
            else:
                bullets = extract_bullets(sec_text)
                if bullets:
                    data["phases"].append({"name": "Lesson Procedure", "content": bullets[:8]})
            break
    for k in sections:
        if 'different' in k or 'scaffold' in k:
            sec_text = sections[k]
            struggling, advanced, ell = [], [], []
            # Try to parse sub-sections by label
            current_group = None
            for line in sec_text.split('\n'):
                ll = line.strip().lower()
                if ll.startswith(('**struggling', '**support', 'struggling', 'support')):
                    current_group = 'struggling'
                elif ll.startswith(('**advanced', '**enrich', '**extend', 'advanced', 'enrich')):
                    current_group = 'advanced'
                elif ll.startswith(('**ell', '**english language', '**language', 'ell', 'english language')):
                    current_group = 'ell'
                elif line.strip().startswith(('- ', '* ', '• ')) or re.match(r'^\d+\.\s', line.strip()):
                    item = re.sub(r'^[-*•\d.]+\s*', '', line.strip()).strip()
                    if item:
                        if current_group == 'struggling':   struggling.append(item)
                        elif current_group == 'advanced':   advanced.append(item)
                        elif current_group == 'ell':        ell.append(item)
            # Fallback: split all bullets into thirds
            if not struggling and not advanced and not ell:
                bullets = extract_bullets(sec_text)
                third = max(1, len(bullets) // 3)
                struggling = bullets[:third]
                advanced   = bullets[third:2*third]
                ell        = bullets[2*third:]
            data["differentiation"]["struggling"] = struggling[:5]
            data["differentiation"]["advanced"]   = advanced[:5]
            data["differentiation"]["ell"]        = ell[:5]
            break
    for k in sections:
        if 'assessment' in k and 'authentic' not in k:
            bullets = extract_bullets(sections[k])
            half = max(1, len(bullets) // 2)
            data["assessment"]["formative"] = bullets[:half]
            data["assessment"]["summative"] = bullets[half:]
            break
    for k in sections:
        if 'reflect' in k:
            data["reflection"] = extract_bullets(sections[k])[:6]; break

    return data


# ── Main Entry Point ──────────────────────────────────────────

def build_pptx(lesson_md):
    """Build a branded .pptx from lesson markdown. Returns BytesIO."""
    parsed = parse_lesson_markdown(lesson_md)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 — Title
    _slide_title(prs, parsed)

    # 2 — Objectives (numbered, with overflow)
    if parsed["objectives"]:
        _slides_bullets(prs, "Learning Objectives",
                        parsed["objectives"],
                        subtitle="Students Will Be Able To (SWBAT)",
                        max_per=5, bullet_char="▸")

    # 3 — Materials (2-column)
    if parsed["materials"]:
        _slide_materials_2col(prs, parsed)

    # 4 — 21st Century Skills (chips)
    if parsed["skills"]:
        _slide_skills_chips(prs, parsed)

    # 5–N — Lesson phases (split teacher/student if present, overflow if long)
    for phase in parsed["phases"]:
        _slide_phase_split(prs, phase)

    if not parsed["phases"]:
        slide, ct = _content_slide(prs, "Lesson Procedure")
        _add_textbox(slide, "See lesson plan for detailed procedure.",
                     CONTENT_LEFT, ct + Inches(0.3), CONTENT_W, Inches(0.5),
                     font_size=16, color=GREY)

    # Differentiation
    _slide_differentiation(prs, parsed)

    # Assessment
    _slide_assessment(prs, parsed)

    # Reflection
    _slide_reflection(prs, parsed)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
